#
# Copyright 2021 Open Raven Inc. and the Mockingbird project authors
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     https://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
#


from typing import final

from pptx import Presentation

from mockingbird.__base import __BaseDocument
from .__base import __BaseUnstructuredDataType


class PPTXDocument(__BaseDocument):

    @final
    def __init__(self):
        super().__init__(extension="pptx")

        # Create a list of docx formats we're going to export.
        self._docx_styles = []
        active_styles = self._configurable_dict["unstructured_data"]["pptx_document"]["active_styles"]

        if active_styles["paragraph_style"]:
            self._docx_styles.append(_PPTXParagraphStyle)

        if active_styles["bullet_point_style"]:
            self._docx_styles.append(_PPTXBulletPointStyle)

    @final
    def save(self, save_path: str) -> None:

        for style in self._docx_styles:
            instantiated_style = style()
            instantiated_style.clone_sensitive_data(other=self)
            instantiated_style.save(save_path=save_path)
            self._meta_data_object.add_other_meta_data(instantiated_style._meta_data_object)


class _PPTXParagraphStyle(__BaseUnstructuredDataType):
    """
    Writes a simple paragraph containing sensitive-soup.
    """

    def __init__(self):
        super().__init__(extension="pptx")

    @final
    def save(self, save_path: str) -> None:
        save_file = self.setup_save_file(save_path=save_path, extension=self.extension)

        sensitive_soup = self._get_sensitive_soup()

        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "A simple title / subtitle slide"
        subtitle.text = sensitive_soup

        prs.save(save_file)
        self._log_save(save_file)


class _PPTXBulletPointStyle(__BaseUnstructuredDataType):
    """
    Writes a simple powerpoint with sensitive-soup in a bullet point.
    """

    def __init__(self):
        super().__init__(extension="pptx")

    @final
    def save(self, save_path: str) -> None:
        """

        """
        save_file = self.setup_save_file(save_path=save_path, extension=self.extension)
        enumerated_groups = self._get_enumerated_style()

        prs = Presentation()

        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "An executive meeting about productivity"
        subtitle.text = "The executive put sensitive information in one of his productivity bullet points to feel " \
                        "productive "

        bullet_slide_layout = prs.slide_layouts[1]

        for group in enumerated_groups:
            key, enumerated_items = group

            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = "Productivity is up 10%"

            tf = body_shape.text_frame
            tf.text = key

            for x in range(len(enumerated_items)):
                p = tf.add_paragraph()
                p.text = enumerated_items[x]
                p.level = 1

        prs.save(save_file)
        self._log_save(save_file)
